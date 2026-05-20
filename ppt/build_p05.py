#!/usr/bin/env python3
"""Rebuild P05-测试与数据质量.pptx with dark theme + full content."""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = 9144000
SLIDE_H = 6858000

BG      = RGBColor(0x1a, 0x1b, 0x26)
TITLE   = RGBColor(0x00, 0x9d, 0xff)
SUB     = RGBColor(0x88, 0xbb, 0xdd)
WHITE   = RGBColor(0xff, 0xff, 0xff)
BODY    = RGBColor(0xcc, 0xcc, 0xdd)
ACCENT  = RGBColor(0x00, 0xcc, 0x99)
ORANGE  = RGBColor(0xff, 0x99, 0x44)
MUTED   = RGBColor(0x66, 0x77, 0x88)
GREEN   = RGBColor(0x00, 0xcc, 0x66)
RED     = RGBColor(0xff, 0x55, 0x55)
YELLOW  = RGBColor(0xff, 0xcc, 0x44)
CARD_BG = RGBColor(0x22, 0x24, 0x32)
CODE_BG = RGBColor(0x2a, 0x2b, 0x3a)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_box(slide, left, top, w, h, text, size=18, color=BODY, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, w, h, items, size=14, color=BODY, title=None, title_color=TITLE):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Arial"
        first = False
    for item in items:
        txt, lvl = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.text = f"{'  ' * lvl}◆ {txt}"
        p.font.size = Pt(size - lvl * 2)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.level = lvl
        first = False
    return txBox

def add_card(slide, left, top, w, h, fill=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_code_block(slide, left, top, w, h, code_text, size=11, color=ACCENT):
    add_card(slide, left, top, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(Emu(left + 100000), Emu(top + 80000), Emu(w - 200000), Emu(h - 160000))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Courier New"

def add_header(slide, text, color=TITLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x15, 0x20)
    shape.line.fill.background()
    add_box(slide, 400000, 120000, 8300000, 600000, text, size=26, color=color, bold=True)

def add_footer(slide, text="P05 - 测试与数据质量 | Nick"):
    add_box(slide, 400000, SLIDE_H - 400000, 8000000, 300000, text, size=10, color=MUTED)

def add_divider(slide, y=3200000):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 400000, Emu(y), Emu(8300000), Emu(3000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
    shape.line.fill.background()


# ═══ SLIDE 1: Title ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1800000, 7300000, 1200000, "P05", size=72, color=TITLE, bold=True, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 2800000, 7300000, 600000, "测试与数据质量", size=36, color=WHITE, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 3400000, 7300000, 400000, "dbt 从入门到精通", size=22, color=SUB, align=PP_ALIGN.LEFT)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 900000, Emu(3900000), Emu(2000000), Emu(3000))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_box(slide, 900000, 4100000, 7300000, 400000, "Nick  |  github.com/raycdut/dbt-learn", size=16, color=MUTED)
add_box(slide, 900000, 4400000, 7300000, 400000, "《高级工程师视角下的 dbt 通关指南》系列", size=14, color=MUTED)


# ═══ SLIDE 2: 为什么重要 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "📋  为什么数据质量重要？")

add_bullets(slide, 500000, 1100000, 8000000, 2500000, [
    "BI 报表数据错误 → 业务决策错误",
    "ML 模型输入错误 → 预测偏差",
    '花 2 小时排查"是不是 dbt 跑错了"',
    "最后发现是源头脏数据",
], size=16, color=BODY)

add_card(slide, 500000, 3800000, 8100000, 1800000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 3900000, 7800000, 400000, "💡 dbt 的解法", size=20, color=ACCENT, bold=True)
add_box(slide, 700000, 4350000, 7800000, 1100000,
    "把数据质量检查写进代码，每次 dbt run 自动验证\n"
    "而不是靠人工定期跑脚本查数据",
    size=15, color=BODY)
add_footer(slide)


# ═══ SLIDE 3: 测试分类 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "1️⃣  测试分类")

tree_text = """dbt test
├── 内置测试 (out-of-the-box)
│   ├── unique        唯一性
│   ├── not_null      非空
│   ├── accepted_values  枚举值
│   └── relationships    外键完整性
│
├── 通用测试 (dbt_utils 包)
│   ├── unique_combination_of_columns
│   ├── expression_is_true
│   ├── recency
│   └── accepted_range
│
├── 自定义测试
│   ├── singular (单个 SQL 文件)
│   └── generic (宏定义，可复用)"""

add_code_block(slide, 400000, 1100000, 5000000, 5000000, tree_text, size=11, color=BODY)

add_card(slide, 5600000, 1100000, 3200000, 3000000, CARD_BG)
add_bullets(slide, 5750000, 1200000, 2900000, 2800000, [
    "内置测试：零配置，YAML 声明即用",
    "dbt_utils：需要 dbt deps 安装",
    "singular test：手写 SQL 查坏数据",
    "generic test：写一次到处复用",
], size=12, color=BODY, title="三类测试", title_color=ORANGE)

add_card(slide, 5600000, 4300000, 3200000, 1800000, CARD_BG)
add_box(slide, 5750000, 4400000, 2900000, 1500000,
    "⚡ 核心规则\n\n"
    "测试返回 0 行 → PASS\n"
    "测试返回 ≥1 行 → FAIL",
    size=14, color=YELLOW, bold=False)
add_footer(slide)


# ═══ SLIDE 4: 内置测试详解 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "2️⃣  内置测试详解")

# 4 cards in 2x2
tests = [
    ("unique", "检查列值是否唯一", """SELECT id, COUNT(*)
FROM {{ ref('model') }}
GROUP BY id
HAVING COUNT(*) > 1"""),
    ("not_null", "检查是否有 NULL 值", """SELECT id
FROM {{ ref('model') }}
WHERE id IS NULL"""),
    ("accepted_values", "检查值在允许范围内", """SELECT DISTINCT status
FROM {{ ref('model') }}
WHERE status NOT IN
  ('a', 'b', 'c')"""),
    ("relationships", "检查外键完整性", """SELECT fk
FROM fact f
WHERE fk NOT IN
  (SELECT pk FROM dim)"""),
]

positions = [
    (400000, 1050000, 4000000, 2500000),
    (4700000, 1050000, 4000000, 2500000),
    (400000, 3750000, 4000000, 2500000),
    (4700000, 3750000, 4000000, 2500000),
]

for (name, desc, sql), (x, y, w, h) in zip(tests, positions):
    add_card(slide, x, y, w, h, CARD_BG)
    add_box(slide, x + 150000, y + 80000, w - 300000, 250000, name, size=16, color=TITLE, bold=True)
    add_box(slide, x + 150000, y + 320000, w - 300000, 200000, desc, size=11, color=SUB)
    add_code_block(slide, x + 100000, y + 550000, w - 200000, h - 700000, sql, size=8, color=ACCENT)

add_footer(slide)


# ═══ SLIDE 5: YAML 声明 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "3️⃣  在 YAML 中声明测试")

add_code_block(slide, 400000, 1100000, 5000000, 4800000, """version: 2
models:
  - name: stg_orders
    columns:
      - name: order_id
        tests:
          - unique
          - not_null
      - name: customer_id
        tests:
          - not_null
      - name: status
        tests:
          - accepted_values:
              values:
                - placed
                - shipped
                - completed
      - name: amount
        tests:
          - not_null
          - dbt_utils.accepted_range:
              min_value: 0""", size=10, color=YELLOW)

add_card(slide, 5600000, 1100000, 3200000, 4800000, CARD_BG)
add_bullets(slide, 5750000, 1200000, 2900000, 4600000, [
    ("内置测试：直接写测试名", 0),
    ("dbt_utils 测试：包名.测试名", 0),
    ("", 0),
    ("测试支持参数：", 0),
    ("accepted_values.values", 1),
    ("relationships.to / field", 1),
    ("dbt_utils 的参数", 1),
    ("min_value / max_value", 1),
    ("", 0),
    ("💡 你项目里已经配了：", 0),
    ("stg_sales_order_header", 1),
    ("  accepted_values: [1..5]", 1),
    ("fct_orders", 1),
    ("  relationships → stg_customer", 1),
    ("  expression_is_true: total_due>0", 1),
], size=12, color=BODY)

add_footer(slide)


# ═══ SLIDE 6: dbt_utils 通用测试 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "4️⃣  dbt_utils — 通用测试包")

add_code_block(slide, 400000, 1000000, 4000000, 1200000, """# packages.yml
packages:
  - package: dbt-labs/dbt_utils
    version: ">=1.0.0"
    
# 安装
dbt deps""", size=11, color=YELLOW)

add_card(slide, 4600000, 1000000, 4100000, 1200000, CARD_BG)
add_box(slide, 4750000, 1100000, 3800000, 1000000,
    "已安装 v1.3.3 ✓\n\n"
    "140+ 通用测试宏，覆盖大多数数据质量场景",
    size=12, color=GREEN)

add_divider(slide, 2400000)

add_bullets(slide, 400000, 2600000, 8300000, 3800000, [
    ("dbt_utils.unique_combination_of_columns:", 0),
    ("复合唯一键检查，比单独 unique 更精确", 1),
    ("", 0),
    ("dbt_utils.expression_is_true:", 0),
    ("自定义表达式断言（已用于 fct_orders.total_due > 0）", 1),
    ("", 0),
    ("dbt_utils.accepted_range:", 0),
    ("值域检查：min_value / max_value", 1),
    ("", 0),
    ("dbt_utils.recency:", 0),
    ("数据新鲜度：最近一次数据更新距今不超过 N 天", 1),
    ("", 0),
    ("dbt_utils.cardinality_equality:", 0),
    ("两层数据量相等检查（源 vs 目标客户数一致）", 1),
], size=13, color=BODY)

add_footer(slide)


# ═══ SLIDE 7: 自定义测试 — Singular ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "5️⃣  自定义测试 — Singular")

add_code_block(slide, 400000, 1050000, 5000000, 2000000, """-- tests/assert_positive_total.sql
SELECT
    sales_order_id,
    total_due
FROM {{ ref('fct_orders') }}
WHERE total_due < 0""", size=12, color=ACCENT)

add_code_block(slide, 400000, 3300000, 5000000, 2000000, """-- tests/assert_valid_order_dates.sql
SELECT
    sales_order_id,
    order_date
FROM {{ ref('fct_orders') }}
WHERE order_date > CURRENT_DATE""", size=12, color=ACCENT)

add_card(slide, 5600000, 1050000, 3200000, 4300000, CARD_BG)
add_bullets(slide, 5750000, 1150000, 2900000, 4100000, [
    ("一个 .sql 文件 = 一个测试", 0),
    ("文件放在 tests/ 目录下", 0),
    ("dbt 自动发现所有 .sql 文件", 0),
    ("", 0),
    ("写法规则：", 0),
    ('SELECT 查出"不该存在"的数据', 1),
    ("返回 0 行 → PASS ✓", 1),
    ("返回 ≥1 行 → FAIL ✗", 1),
    ("", 0),
    ("💡 你项目里已有：", 0),
    ("tests/assert_positive_total.sql", 1),
    ("tests/assert_valid_order_dates.sql", 1),
    ("dbt build 自动运行 ✓", 1),
], size=12, color=BODY, title="📝 Singular Test", title_color=ORANGE)

add_footer(slide)


# ═══ SLIDE 8: 自定义测试 — Generic ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "6️⃣  自定义测试 — Generic")

add_code_block(slide, 400000, 1050000, 5000000, 2200000, """-- tests/generic/test_is_positive.sql
{% test is_positive(model, column_name) %}

SELECT
    {{ column_name }}
FROM {{ model }}
WHERE {{ column_name }} <= 0

{% endtest %}""", size=12, color=ACCENT)

add_code_block(slide, 400000, 3500000, 5000000, 1100000, """# 使用：在 YAML 中引用
columns:
  - name: amount
    tests:
      - is_positive""", size=11, color=YELLOW)

add_card(slide, 5600000, 1050000, 3200000, 3500000, CARD_BG)
add_bullets(slide, 5750000, 1150000, 2900000, 3300000, [
    ("宏模板 + YAML 引用", 0),
    ("写一次，到处复用", 0),
    ("", 0),
    ("两个固定参数：", 0),
    ("model → 被测试的表", 1),
    ("column_name → 被测试的列", 1),
    ("", 0),
    ("位置：tests/generic/ 下", 0),
    ("dbt 自动扫描加载", 0),
    ("", 0),
    ("Singular vs Generic：", 0),
    ("Singular = 一次性断言", 1),
    ("Generic = 可复用 + YAML 配", 1),
], size=12, color=BODY, title="🔁 Generic Test", title_color=ACCENT)

add_footer(slide)


# ═══ SLIDE 9: 运行测试 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "7️⃣  运行测试")

add_code_block(slide, 400000, 1050000, 4200000, 3200000, """# 跑全部测试
dbt test

# 跑指定 model 的测试
dbt test --select fct_orders

# 跑 dbt build（run + test 一步）
dbt build

# 存储失败记录
dbt test --store-failures

# 只跑带标签的测试
dbt test --select tag:critical""", size=11, color=ACCENT)

add_card(slide, 4800000, 1050000, 3900000, 3200000, CARD_BG)
add_box(slide, 4950000, 1150000, 3600000, 300000, "⚡ store_failures", size=16, color=ORANGE, bold=True)
add_box(slide, 4950000, 1500000, 3600000, 2500000,
    "默认：测试失败只报错\n"
    "--store-failures 后：\n"
    "  把坏数据写进表\n"
    "  schema.dbt_test__audit\n\n"
    "方便排查：\n"
    "  SELECT * FROM audit\n"
    "  WHERE test_name = 'unique'\n\n"
    "全局开启：\n"
    "  dbt_project.yml 中配\n"
    "  +store_failures: true",
    size=12, color=BODY)

add_card(slide, 400000, 4500000, 8300000, 1800000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 600000, 4600000, 7900000, 1500000,
    "📊 你的项目跑完结果：\n"
    "132 steps - 11 models, 90 seeds, 31 data tests\n"
    "全部 PASS ✓\n"
    "包括 2 个 singular test + 1 个 dbt_utils.expression_is_true + 1 个 relationships",
    size=13, color=GREEN)

add_footer(slide)


# ═══ SLIDE 10: 测试策略 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "8️⃣  测试策略建议")

# Table header
cols = [400000, 1300000, 3800000, 6300000]
col_w = [900000, 2400000, 2400000, 2600000]
header_texts = ["层", "必测", "建议测", "可选"]
for i, (x, h) in enumerate(zip(cols, header_texts)):
    add_box(slide, x, 1050000, col_w[i], 300000, h, size=12, color=TITLE, bold=True)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 400000, Emu(1400000), Emu(8300000), Emu(2000))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
shape.line.fill.background()

rows = [
    ("Staging", "not_null, unique (主键)", "accepted_values", "-"),
    ("Intermediate", "not_null (关键字段)", "relationships", "cardinality_equality"),
    ("Marts", "not_null, unique (主键)", "relationships, recency", "自定义业务规则"),
]
row_data = [
    ("not_null, unique\n(主键)", "accepted_values", "-"),
    ("not_null\n(关键字段)", "relationships", "cardinality_equality"),
    ("not_null, unique\n(主键)", "relationships\nrecency", "自定义业务规则"),
]
for i, (layer, must, suggest, optional) in enumerate(rows):
    y = 1550000 + i * 800000
    add_card(slide, 400000, y, 8300000, 650000, CARD_BG)
    add_box(slide, 450000, y + 100000, col_w[0], 450000, layer, size=13, color=TITLE, bold=True)
    add_box(slide, cols[1] + 50000, y + 100000, col_w[1], 450000, must, size=11, color=ORANGE)
    add_box(slide, cols[2] + 50000, y + 100000, col_w[2], 450000, suggest, size=11, color=BODY)
    add_box(slide, cols[3] + 50000, y + 100000, col_w[3], 450000, optional, size=11, color=MUTED)

add_divider(slide, 4100000)

add_box(slide, 400000, 4300000, 4000000, 300000, "❌ 不要测什么", size=16, color=WHITE, bold=True)
add_bullets(slide, 400000, 4650000, 4000000, 1800000, [
    "非关键字段的 not_null（允许 NULL）",
    "内部逻辑字段的唯一性（不是业务主键）",
    "跨大表的关系完整性（改抽样）",
], size=12, color=BODY)

add_box(slide, 4600000, 4300000, 4100000, 300000, "⚠️ 测试失败怎么办", size=16, color=WHITE, bold=True)
add_bullets(slide, 4600000, 4650000, 4100000, 1800000, [
    "看日志：哪个 model / 字段",
    "检查源数据是否脏",
    "数据修复：UPDATE / 重跑 seed",
    "加代码防御：stg 层过滤异常",
    "重新 dbt build",
], size=12, color=BODY)

add_footer(slide)


# ═══ SLIDE 11: 同声传译 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "🛠️  同声传译 — 工程师视角")

concepts = [
    ("unique 测试", "dropDuplicates count 对比", "检查列值是否唯一"),
    ("not_null", "filter(col.isNull()).count()", "检查是否有空值"),
    ("accepted_values", "WHERE NOT IN 条件", "值必须在允许范围内"),
    ("relationships", "事实表外键在维度表匹配", "外键约束检查"),
    ("singular test", "一条 SQL 脚本", "自定义数据质量规则"),
    ("generic test", "可复用宏模板", "写一次到处复用"),
    ("store_failures", "_bad_records 分区", "坏数据存下来分析"),
]

y_start = 1050000
for i, (dbt_concept, spark_concept, plain) in enumerate(concepts):
    y = y_start + i * 730000
    add_card(slide, 400000, y, 2500000, 580000, CARD_BG)
    add_box(slide, 450000, y + 50000, 2400000, 160000, dbt_concept, size=12, color=TITLE, bold=True)
    add_box(slide, 450000, y + 220000, 2400000, 300000, spark_concept, size=10, color=BODY)

    add_card(slide, 3000000, y, 2700000, 580000, CARD_BG)
    add_box(slide, 3050000, y + 50000, 2600000, 160000, "你熟悉的场景", size=10, color=ORANGE, bold=True)
    add_box(slide, 3050000, y + 220000, 2600000, 300000, spark_concept, size=10, color=BODY)

    add_card(slide, 5800000, y, 3000000, 580000, CARD_BG)
    add_box(slide, 5850000, y + 50000, 2900000, 160000, "大白话", size=10, color=ACCENT, bold=True)
    add_box(slide, 5850000, y + 220000, 2900000, 300000, plain, size=10, color=BODY)

add_footer(slide)


# ═══ SLIDE 12: 直觉陷阱 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "⚠️  工程师直觉陷阱")

add_card(slide, 500000, 1200000, 8100000, 1200000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 700000, 1300000, 7800000, 1000000,
    "🔥 直觉陷阱\n"
    "\"测试应该在数据摄入时做，不是转换时\"",
    size=16, color=RED)

add_card(slide, 500000, 2600000, 8100000, 1500000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 2700000, 7800000, 1300000,
    "✅ dbt 的方式\n"
    "两个阶段都要做。摄入时做 schema 校验，转换时做业务逻辑校验。\n"
    "比如订单金额不可能为负数——摄入时看不出来，转换层容易检测。\n"
    "类比：Bronze 层做格式检查，Silver/Gold 层做业务合理性检查。",
    size=14, color=ACCENT)

add_card(slide, 500000, 4300000, 8100000, 1200000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 700000, 4400000, 7800000, 1000000,
    "💡 实战映射\n"
    "你在药企架构里讨论的 cross_source_consistency 测试——\n"
    "不同数据源的同一个客户信息是否一致。\n"
    "用 dbt 的 relationships test + dbt_expectations 包可以直接写。",
    size=13, color=SUB)
add_footer(slide)


# ═══ SLIDE 13: 下集预告 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1500000, 7300000, 800000, "✅  本集完成", size=40, color=GREEN, bold=True)

add_box(slide, 900000, 2500000, 7300000, 300000, "你已学会：", size=16, color=WHITE, bold=True)
add_bullets(slide, 900000, 2800000, 7300000, 1500000, [
    "三类测试：内置 / dbt_utils / 自定义",
    "YAML 声明测试：unique, not_null, accepted_values, relationships",
    "dbt_utils 通用测试包安装与使用",
    "Singular + Generic 自定义测试",
    "运行策略：dbt test / build / store-failures",
    "分层测试策略：stg → int → marts",
], size=14, color=BODY)

add_divider(slide, 4500000)

add_box(slide, 900000, 4700000, 7300000, 400000, "下一集：P06 — 自动文档 + 血缘图", size=18, color=TITLE, bold=True)

add_box(slide, 900000, 5400000, 7300000, 400000, "配套代码：github.com/raycdut/dbt-learn", size=14, color=MUTED)
add_box(slide, 900000, 5700000, 7300000, 400000, "我是 Nick，下期见 👋", size=14, color=SUB)


# ── Save ──
output_path = "/Users/chendong/Projects/dbt-learn/ppt/P05-测试.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
