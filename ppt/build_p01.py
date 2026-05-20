#!/usr/bin/env python3
"""Rebuild P01-环境搭建.pptx with dark theme and full content."""

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──
SLIDE_W = 9144000   # 10 inches
SLIDE_H = 6858000   # 7.5 inches

# Color palette (dark theme)
BG      = RGBColor(0x1a, 0x1b, 0x26)
TITLE   = RGBColor(0x00, 0x9d, 0xff)
SUB     = RGBColor(0x88, 0xbb, 0xdd)
WHITE   = RGBColor(0xff, 0xff, 0xff)
BODY    = RGBColor(0xcc, 0xcc, 0xdd)
ACCENT  = RGBColor(0x00, 0xcc, 0x99)
ORANGE  = RGBColor(0xff, 0x99, 0x44)
MUTED   = RGBColor(0x66, 0x77, 0x88)
GREEN   = RGBColor(0x00, 0xcc, 0x66)
RED     = RGBColor(0xff, 0x55, 0x55)
YELLOW  = RGBColor(0xff, 0xcc, 0x44)
CARD_BG = RGBColor(0x22, 0x24, 0x32)
CODE_BG = RGBColor(0x2a, 0x2b, 0x3a)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank


# ── Helpers ──
def add_bg(slide, color=BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_box(slide, left, top, w, h, text, size=18, color=BODY, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, w, h, items, size=14, color=BODY, title=None, title_color=TITLE):
    """items: list of (text, indent_level=0) or plain strings"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Arial"
        first = False
    for item in items:
        if isinstance(item, str):
            txt, lvl = item, 0
        else:
            txt, lvl = item
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.text = f"{'  ' * lvl}◆ {txt}"
        p.font.size = Pt(size - lvl * 2)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.level = lvl
        first = False
    return txBox

def add_card(slide, left, top, w, h, fill=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_code_block(slide, left, top, w, h, code_text, size=11, color=ACCENT):
    """Monospace code display with card background"""
    shape = add_card(slide, left, top, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(Emu(left + 100000), Emu(top + 80000), Emu(w - 200000), Emu(h - 160000))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Courier New"

def add_header(slide, text, color=TITLE):
    """Dark header banner at the top"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x15, 0x20)
    shape.line.fill.background()
    add_box(slide, 400000, 120000, 8300000, 600000, text, size=26, color=color, bold=True)

def add_footer(slide, text="P01 - 环境搭建 + 第一个项目 | Nick"):
    add_box(slide, 400000, SLIDE_H - 400000, 8000000, 300000, text, size=10, color=MUTED)

def add_divider(slide, y=3200000):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 400000, Emu(y), Emu(8300000), Emu(3000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
    shape.line.fill.background()


# ════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Decorative top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

# Title
add_box(slide, 900000, 1800000, 7300000, 1200000, "P01", size=72, color=TITLE, bold=True, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 2800000, 7300000, 600000, "环境搭建 + 第一个项目", size=36, color=WHITE, bold=False, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 3400000, 7300000, 400000, "dbt 从入门到精通", size=22, color=SUB, align=PP_ALIGN.LEFT)

# Divider line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 900000, Emu(3900000), Emu(2000000), Emu(3000))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Info
add_box(slide, 900000, 4100000, 7300000, 400000, "Nick  |  github.com/raycdut/dbt-learn", size=16, color=MUTED)
add_box(slide, 900000, 4400000, 7300000, 400000, "《高级工程师视角下的 dbt 通关指南》系列", size=14, color=MUTED)


# ════════════════════════════════════════════
# SLIDE 2: 今天做什么
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "📋  今天做什么")

add_bullets(slide, 500000, 1100000, 8000000, 5000000, [
    "搭建 Python 虚拟环境 + 安装 dbt-core + dbt-duckdb",
    "dbt init 创建项目，认识目录结构",
    "配置 ~/.dbt/profiles.yml 连接 DuckDB",
    "写第一个 model：SELECT 1 AS id",
    "dbt run 执行 + dbt debug 调试",
    "duckdb 命令行验证结果",
    "了解 dbt 执行流程",
], size=16, color=BODY)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 3: 虚拟环境搭建
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "1️⃣  Python 虚拟环境")

# Left pane: explanation
add_bullets(slide, 500000, 1100000, 3800000, 2500000, [
    "dbt 是 Python 包，需要独立环境",
    "防止依赖冲突（dbt 库版本要求严格）",
    "推荐用 uv（比 pip 快 10 倍）",
    "也可以用 Python 自带 venv",
], size=14, color=BODY)

# Right pane: code
add_code_block(slide, 4800000, 1100000, 4000000, 2300000, """# 方法一：使用 uv（推荐）
brew install uv
uv venv -p 3.12 ~/.venvs/dbt
source ~/.venvs/dbt/bin/activate

# 方法二：原生 venv
python3.12 -m venv ~/.venvs/dbt
source ~/.venvs/dbt/bin/activate

pip install dbt-core dbt-duckdb
dbt --version""", size=10, color=ACCENT)

# Warning block
shape = add_card(slide, 500000, 3800000, 8100000, 1000000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 700000, 3900000, 7800000, 800000,
    "⚠️  Python 3.14 不兼容！dbt-core 1.11.x 在 Python 3.14 上会报错。\n请使用 Python 3.12 或 3.13。  检查版本：python3 --version", size=12, color=RED)

# DuckDB note
add_box(slide, 500000, 5000000, 8100000, 600000,
    "💡  为什么选 DuckDB？本地文件数据库，零配置，学完切 Snowflake 只需换一行连接串", size=12, color=ACCENT)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 4: 项目初始化
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "2️⃣  dbt init — 创建项目")

# Code
add_code_block(slide, 500000, 1000000, 3800000, 1200000, """dbt init dbt_learn
cd dbt_learn""", size=12, color=ACCENT)

# Directory structure
dir_tree = """dbt_learn/
├── .gitignore
├── README.md
├── analyses/          ← 临时分析
├── dbt_project.yml    ← 核心配置 ★
├── macros/            ← Jinja 宏
├── models/            ← 模型 SQL  ★★★
│   └── example/
├── seeds/             ← CSV 数据
├── snapshots/         ← SCD 快照
└── tests/             ← 数据测试"""

add_code_block(slide, 500000, 2400000, 4200000, 3800000, dir_tree, size=10, color=BODY)

# Explanation
add_bullets(slide, 5000000, 1100000, 3800000, 4000000, [
    ("80% 的时间在 models/ 里写 SQL", 0),
    ("dbt_project.yml 控制所有配置", 0),
    ("seeds/ 放 CSV 自动建表", 0),
    ("约定优于配置 → 目录结构固定", 0),
    ("与 Databricks 三层架构对应：", 0),
    ("staging / intermediate / marts", 1),
    ("→ bronze / silver / gold", 1),
], size=13, color=BODY)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 5: 配置数据库连接
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "3️⃣  配置数据库连接 — profiles.yml")

# Left: profiles.yml
add_code_block(slide, 400000, 1100000, 4200000, 2400000, """~/.dbt/profiles.yml

dbt_learn:
  outputs:
    dev:
      type: duckdb
      path: dev.duckdb
      threads: 4
  target: dev""", size=11, color=YELLOW)

# Right: explanation
add_bullets(slide, 4900000, 1100000, 4000000, 2500000, [
    "dbt 本身不存数据",
    "它生成 SQL 去数据库执行",
    "profiles.yml 告诉 dbt 连哪个库",
    "DuckDB = 本地文件，无需服务",
    "target: dev → 可加 prod",
], size=13, color=BODY)

# dbt debug
add_box(slide, 400000, 3700000, 8300000, 400000, "验证连接：dbt debug  →  All checks passed! ✓", size=14, color=GREEN)

# Bottom: dbt_project.yml key
add_code_block(slide, 400000, 4300000, 8300000, 2000000, """dbt_project.yml 要点：
name: 'dbt_learn'        # 项目名
profile: 'dbt_learn'     # 对应 profiles.yml 的 key
model-paths: ['models']  # 模型目录
models:
  dbt_learn:
    +materialized: view  # 默认物化方式：view""", size=10, color=ACCENT)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 6: 第一个 model
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "4️⃣  写第一个 model")

# Step 1
add_box(slide, 500000, 1100000, 4000000, 300000, "① 清掉示例文件", size=16, color=WHITE, bold=True)
add_code_block(slide, 500000, 1400000, 4000000, 500000, "rm -rf models/example", size=12, color=ACCENT)

# Step 2
add_box(slide, 500000, 2100000, 4000000, 300000, "② 创建 models/my_first_model.sql", size=16, color=WHITE, bold=True)
add_code_block(slide, 500000, 2400000, 4000000, 600000, """SELECT 1 AS id,
  'Hello dbt!' AS message""", size=12, color=ACCENT)

# Step 3
add_box(slide, 500000, 3200000, 4000000, 300000, "③ 执行", size=16, color=WHITE, bold=True)
add_code_block(slide, 500000, 3500000, 4000000, 800000, """dbt run

1 of 1 OK created sql view model
  main.my_first_model [OK in 0.02s]""", size=12, color=GREEN)

# Right: key insight
add_card(slide, 4800000, 1100000, 4000000, 3200000, CARD_BG)
add_bullets(slide, 4950000, 1200000, 3700000, 3000000, [
    ("💡 一个 model = 一个 .sql 文件", 0),
    ("里面就放一条 SELECT 语句", 0),
    ("不需要 INSERT / CREATE TABLE", 0),
    ("dbt 自动用 SELECT 建 view/table", 0),
    ("", 0),
    ("默认物化方式：view", 0),
    ("每次 dbt run 重建 view", 0),
    ("改 +materialized: table", 0),
    ("就变成 DROP+CREATE TABLE", 0),
], size=13, color=BODY, title="✨ model 哲学", title_color=ORANGE)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 7: 验证
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "5️⃣  验证结果")

add_code_block(slide, 400000, 1200000, 8300000, 1200000, """# 用 DuckDB CLI 直接查
duckdb dev.duckdb -c "SELECT * FROM my_first_model"

id | message
1  | Hello dbt!""", size=13, color=ACCENT)

add_divider(slide, 2800000)

add_box(slide, 400000, 3100000, 8300000, 400000, "dbt debug  — 诊断连接和配置问题", size=14, color=ORANGE)
add_box(slide, 400000, 3500000, 8300000, 400000, "dbt compile  — 编译生成完整 SQL（不执行），查看 dbt 到底跑了什么", size=14, color=BODY)

add_divider(slide, 4100000)

# Execution flow explanation
add_bullets(slide, 400000, 4300000, 8300000, 2300000, [
    "dbt run 执行流程：",
    ("读 dbt_project.yml → 找 models/ → 解析依赖 → 编译 SQL → 执行", 1),
    "相当于 spark-submit 的简化版：解析 DAG → 按序执行",
    "DuckDB 数据存在 dev.duckdb 文件里，打开就能看",
], size=13, color=BODY)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 8: 操作清单
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "✅  操作清单")

add_bullets(slide, 500000, 1100000, 8000000, 4500000, [
    "创建 Python 虚拟环境（uv 或 venv）",
    "pip install dbt-core dbt-duckdb",
    "dbt --version 验证安装",
    "dbt init dbt_learn 创建项目",
    "配置 ~/.dbt/profiles.yml（DuckDB 连接）",
    "写第一个 model：SELECT 1 AS id",
    "dbt run 执行",
    "dbt compile 查看编译后的 SQL",
    "dbt debug 诊断问题",
    "duckdb dev.duckdb 验证数据",
], size=14, color=BODY)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 9: 同声传译（工程师视角）
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "🛠️  同声传译 — 工程师视角")

# Table-like layout using cards
concepts = [
    ("dbt-core / dbt Cloud", "Databricks Community vs Workspace", "免费本地 vs SaaS 付费"),
    ("adapter", "Spark JDBC connector", "连不同数据库的插件"),
    ("profiles.yml", "spark-defaults.conf", "数据库连接配置"),
    ("dbt_project.yml", "build.sbt / pom.xml", "项目元信息配置"),
    ("dbt init", "mvn archetype:generate", "一键脚手架生成"),
    ("model (.sql)", "DataFrame transformation", "一个文件 = 一个转换步骤"),
]

y_start = 1100000
for i, (dbt, spark, plain) in enumerate(concepts):
    y = y_start + i * 700000
    add_card(slide, 400000, y, 2500000, 550000, CARD_BG)
    add_box(slide, 450000, y + 50000, 2400000, 150000, dbt, size=13, color=TITLE, bold=True)
    add_box(slide, 450000, y + 200000, 2400000, 300000, spark, size=11, color=BODY)
    
    add_card(slide, 3000000, y, 2700000, 550000, CARD_BG)
    add_box(slide, 3050000, y + 50000, 2600000, 150000, "你熟悉的场景", size=11, color=ORANGE, bold=True)
    add_box(slide, 3050000, y + 200000, 2600000, 300000, spark, size=11, color=BODY)
    
    add_card(slide, 5800000, y, 3000000, 550000, CARD_BG)
    add_box(slide, 5850000, y + 50000, 2900000, 150000, "大白话", size=11, color=ACCENT, bold=True)
    add_box(slide, 5850000, y + 200000, 2900000, 300000, plain, size=11, color=BODY)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 10: 直觉陷阱
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "⚠️  工程师直觉陷阱")

add_card(slide, 500000, 1200000, 8100000, 1200000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 700000, 1300000, 7800000, 1000000,
    "🔥 直觉陷阱\n"
    "\"dbt 初始化太简单了，目录模板会不会不够灵活？\"",
    size=16, color=RED)

add_card(slide, 500000, 2600000, 8100000, 1500000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 2700000, 7800000, 1300000,
    "✅ dbt 的哲学：约定优于配置\n"
    "• 目录结构是固定的：staging / intermediate / marts\n"
    "• 跟 Databricks 的 bronze / silver / gold 完全对应\n"
    "• 不要自己改目录，按约定来，别人秒懂",
    size=14, color=ACCENT)

add_card(slide, 500000, 4300000, 8100000, 1200000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 700000, 4400000, 7800000, 1000000,
    "💡 实战映射\n"
    "你在药企架构里的 Raw→Conformed→Enriched→Curated 三层思路，和 dbt 的目录设计完全一致。\n"
    "区别只是：之前在 Databricks 手工建文件夹，dbt 帮你强制了这个规范。",
    size=13, color=SUB)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 11: 常见问题
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "❓  常见问题")

problems = [
    ("dbt 命令找不到？", "没有激活虚拟环境：source ~/.venvs/dbt/bin/activate"),
    ("profiles.yml 报错？", "dbt debug 诊断连接 → 检查 type/path 是否匹配"),
    ("model 没建出来？", "dbt compile 编译查看生成的 SQL → 检查语法错误"),
    ("DuckDB 文件在哪？", "当前目录下 dev.duckdb 文件，可用 SQLite 浏览器打开"),
]

for i, (q, a) in enumerate(problems):
    y = 1200000 + i * 1200000
    add_box(slide, 600000, y, 8000000, 300000, f"Q：{q}", size=15, color=YELLOW, bold=True)
    add_box(slide, 800000, y + 400000, 7800000, 500000, f"A：{a}", size=13, color=BODY)
    if i < len(problems) - 1:
        add_divider(slide, y + 950000)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 12: 下集预告 + 感谢
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Decorative top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1500000, 7300000, 800000, "✅  本集完成", size=40, color=GREEN, bold=True)

add_box(slide, 900000, 2500000, 7300000, 300000, "你已学会：", size=16, color=WHITE, bold=True)
add_bullets(slide, 900000, 2800000, 7300000, 1500000, [
    "搭建 Python 虚拟环境 + 安装 dbt",
    "dbt init 创建项目，了解目录结构",
    "写 model + dbt run 执行 + 验证",
], size=14, color=BODY)

add_divider(slide, 4500000)

add_box(slide, 900000, 4700000, 7300000, 400000, "下一集：P02 — ref / source / model：dbt 的灵魂三件套", size=18, color=TITLE, bold=True)

add_box(slide, 900000, 5400000, 7300000, 400000, "配套代码：github.com/raycdut/dbt-learn", size=14, color=MUTED)
add_box(slide, 900000, 5700000, 7300000, 400000, "我是 Nick，下期见 👋", size=14, color=SUB)


# ── Save ──
output_path = "/Users/chendong/projects/dbt-learn/ppt/P01-环境搭建.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
