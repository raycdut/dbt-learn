#!/usr/bin/env python3
"""Rebuild P06-文档.pptx with dark theme + full content."""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = 9144000
SLIDE_H = 6858000

BG      = RGBColor(0x1a, 0x1b, 0x26)
TITLE   = RGBColor(0x00, 0x9d, 0xff)
SUB     = RGBColor(0x88, 0xbb, 0xdd)
WHITE   = RGBColor(0xff, 0xff, 0xff)
BODY    = RGBColor(0xcc, 0xcc, 0xdd)
ACCENT  = RGBColor(0x00, 0xcc, 0x99)
ORANGE  = RGBColor(0xff, 0x99, 0x44)
MUTED   = RGBColor(0x66, 0x77, 0x88)
GREEN   = RGBColor(0x00, 0xcc, 0x66)
RED     = RGBColor(0xff, 0x55, 0x55)
YELLOW  = RGBColor(0xff, 0xcc, 0x44)
CARD_BG = RGBColor(0x22, 0x24, 0x32)
CODE_BG = RGBColor(0x2a, 0x2b, 0x3a)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_box(slide, left, top, w, h, text, size=18, color=BODY, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, w, h, items, size=14, color=BODY, title=None, title_color=TITLE):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = title_color
        p.font.bold = True
        first = False
    for item in items:
        txt, lvl = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.text = f"{'  ' * lvl}◆ {txt}"
        p.font.size = Pt(size - lvl * 2)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.level = lvl
        first = False
    return txBox

def add_card(slide, left, top, w, h, fill=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_code_block(slide, left, top, w, h, code_text, size=11, color=ACCENT):
    add_card(slide, left, top, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(Emu(left + 100000), Emu(top + 80000), Emu(w - 200000), Emu(h - 160000))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Courier New"

def add_header(slide, text, color=TITLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x15, 0x20)
    shape.line.fill.background()
    add_box(slide, 400000, 120000, 8300000, 600000, text, size=26, color=color, bold=True)

def add_footer(slide, text="P06 - 自动文档 + 血缘图 | Nick"):
    add_box(slide, 400000, SLIDE_H - 400000, 8000000, 300000, text, size=10, color=MUTED)

def add_divider(slide, y=3200000):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 400000, Emu(y), Emu(8300000), Emu(3000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
    shape.line.fill.background()


# ═══ SLIDE 1: Title ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1800000, 7300000, 1200000, "P06", size=72, color=TITLE, bold=True, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 2800000, 7300000, 600000, "自动文档 + 血缘图", size=36, color=WHITE, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 3400000, 7300000, 400000, "dbt 从入门到精通", size=22, color=SUB, align=PP_ALIGN.LEFT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 900000, Emu(3900000), Emu(2000000), Emu(3000))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
add_box(slide, 900000, 4100000, 7300000, 400000, "Nick  |  github.com/raycdut/dbt-learn", size=16, color=MUTED)
add_box(slide, 900000, 4400000, 7300000, 400000, "《高级工程师视角下的 dbt 通关指南》系列", size=14, color=MUTED)


# ═══ SLIDE 2: 为什么重要 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "📋  为什么自动文档重要？")

add_bullets(slide, 500000, 1100000, 8000000, 2500000, [
    '新同事问：\"这张 fct_orders 表啥意思？\"',
    '你回：\"就是订单…吧？\"',
    '查 10 分钟源码才发现是\"已支付的订单\"',
    '不同表里同名字段的口径不一致',
], size=16, color=BODY)

add_card(slide, 500000, 3800000, 8100000, 1800000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 3900000, 7800000, 400000, "💡 dbt 的解法", size=20, color=ACCENT, bold=True)
add_box(slide, 700000, 4350000, 7800000, 1100000,
    "在 schema.yml 里写描述\n"
    "→ 一键生成文档站（dbt docs generate）\n"
    "→ 改代码时顺手更新，再也不用追着人补文档",
    size=15, color=BODY)
add_footer(slide)


# ═══ SLIDE 3: 写描述 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "1️⃣  写描述 — model + column + source")

add_code_block(slide, 400000, 1050000, 5000000, 3000000, """version: 2
models:
  - name: fct_orders
    description: |-
      订单事实表。
      包含所有已支付订单明细，
      每日增量更新。
      上游：stg_orders
      下游：BI 报表
    columns:
      - name: order_id
        description: "订单唯一标识"
      - name: amount
        description: "金额（含税）"
      - name: status
        description: "placed/shipped/..." """, size=9, color=YELLOW)

add_card(slide, 5600000, 1050000, 3200000, 3000000, CARD_BG)
add_bullets(slide, 5750000, 1150000, 2900000, 2800000, [
    ("三层描述都要写：", 0),
    ("model 级 → 表的业务含义", 1),
    ("column 级 → 字段口径", 1),
    ("source 级 → 数据来源", 1),
    ("", 0),
    ("支持 YAML 多行 | 语法", 0),
    ("支持 Markdown 格式", 0),
    ("跟代码一起 PR 审查", 0),
], size=12, color=BODY)

# 实操提示
add_card(slide, 400000, 4300000, 8300000, 1800000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 600000, 4400000, 7900000, 1500000,
    "📝 你的项目已配置：\n"
    "• 所有 stg/int/marts 层 model 已加 description\n"
    "• fct_orders 用了 docs block（Markdown 富文档）\n"
    "• stg_sales_order_header 也配了 docs block\n"
    "• 文档和代码同 PR，更新即可同步",
    size=12, color=ACCENT)
add_footer(slide)


# ═══ SLIDE 4: 生成文档 + 启动 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "2️⃣  生成文档 & 启动服务")

add_code_block(slide, 400000, 1050000, 5000000, 1600000, """# 1. 生成文档（需先 dbt run）
dbt docs generate

# 2. 启动本地文档站
dbt docs serve

# 浏览器打开 http://localhost:8080""", size=13, color=ACCENT)

add_card(slide, 5600000, 1050000, 3200000, 1600000, CARD_BG)
add_bullets(slide, 5750000, 1150000, 2900000, 1400000, [
    ("dbt docs generate：", 0),
    ("读取 schema.yml → catalog.json", 1),
    ("生成静态 HTML 文档站", 1),
    ("dbt docs serve：", 0),
    ("启动 Flask 服务预览", 1),
], size=12, color=BODY)

add_divider(slide, 2900000)

add_box(slide, 400000, 3100000, 8300000, 300000, "你能在文档站看到什么", size=18, color=WHITE, bold=True)

# 2x2 grid
items = [
    ("📋 项目总览", "所有 model 列表 + 数据源 + 测试统计"),
    ("🔍 单 model 详情", "SQL 代码 + 字段描述 + 属性 + 测试"),
    ("🔗 血缘图", "节点可跳转，展开上下游，缩放查看"),
    ("🏷️ Tag 筛选", "按标签快速定位核心模型"),
]
for i, (title, desc) in enumerate(items):
    x = 400000 + (i % 2) * 4200000
    y = 3550000 + (i // 2) * 1000000
    add_card(slide, x, y, 4000000, 850000, CARD_BG)
    add_box(slide, x + 100000, y + 80000, 3800000, 250000, title, size=14, color=TITLE, bold=True)
    add_box(slide, x + 100000, y + 350000, 3800000, 400000, desc, size=12, color=BODY)

add_footer(slide)


# ═══ SLIDE 5: docs block ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "3️⃣  高级文档 — docs block")

add_code_block(slide, 400000, 1050000, 5000000, 2500000, """-- models/docs.md
{% docs fct_orders_description %}

# 订单事实表

## 业务口径
包含**全部已支付**的订单记录。
已取消、未支付不在此表。

## 刷新策略
每日 08:00 增量刷新。
保留最近 3 年数据。

{% enddocs %}""", size=10, color=ACCENT)

add_code_block(slide, 400000, 3750000, 5000000, 1000000, """# schema.yml 引用
models:
  - name: fct_orders
    description: "{{ doc('fct_orders_description') }}" """, size=10, color=YELLOW)

add_card(slide, 5600000, 1050000, 3200000, 3700000, CARD_BG)
add_bullets(slide, 5750000, 1150000, 2900000, 3500000, [
    ("为什么需要 docs block？", 0),
    ("", 0),
    ("长文档写在 YAML 里：", 0),
    ("→ 格式难看、难编辑", 1),
    ("→ 字符串转义问题", 1),
    ("→ 多人编辑冲突", 1),
    ("", 0),
    ("docs block 解法：", 0),
    ("独立的 .md 文件", 1),
    ("支持完整 Markdown", 1),
    ("{% docs name %}...{% enddocs %}", 1),
    ("YAML 中用 {{ doc('name') }} 引用", 1),
    ("", 0),
    ("💡 你项目已有：", 0),
    ("models/docs.md", 1),
    ("fct_orders_description", 1),
    ("fct_orders_total_due", 1),
], size=10, color=BODY)

add_footer(slide)


# ═══ SLIDE 6: tag + 最佳实践 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "4️⃣  Tag 筛选 + 文档最佳实践")

# Tag section
add_box(slide, 400000, 1050000, 4000000, 300000, "Tag 筛选", size=18, color=WHITE, bold=True)
add_code_block(slide, 400000, 1400000, 4000000, 1200000, """# YAML 中加 tag
models:
  - name: fct_orders
    config:
      tags: ['finance', 'core']

# 文档站按 tag 筛选""", size=11, color=ACCENT)

add_card(slide, 4600000, 1050000, 4100000, 1550000, CARD_BG)
add_box(slide, 4750000, 1150000, 3800000, 1300000,
    "🏷️ 你的项目已配：\n"
    "  fct_orders → finance, core\n"
    "  stg_sales_order_header → core",
    size=13, color=ACCENT)

add_divider(slide, 2900000)

# Best practices
add_box(slide, 400000, 3050000, 4000000, 300000, "✅ 什么该写", size=16, color=WHITE, bold=True)
add_bullets(slide, 400000, 3400000, 4000000, 2200000, [
    "表的业务含义",
    "字段的计算口径",
    "数据来源",
    "刷新频率",
    "数据范围",
    "约束（unique / not_null）",
], size=12, color=GREEN)

add_box(slide, 4600000, 3050000, 4100000, 300000, "❌ 什么不该写", size=16, color=WHITE, bold=True)
add_bullets(slide, 4600000, 3400000, 4100000, 2200000, [
    "技术实现细节（看了 SQL 就懂）",
    "人人都知道（\"id 是标识符\"）",
    "今天写明天忘的废话（\"这个字段很重要\"）",
], size=12, color=RED)

add_card(slide, 400000, 4800000, 8300000, 1500000, CARD_BG)
add_box(slide, 600000, 4900000, 7900000, 1300000,
    "📝 文档模板 — 复制到你的项目填空：\n"
    "description: |\n"
    "  [一句话说明模型的作用]\n"
    "  ## 业务口径\n"
    "  [口径逻辑]\n"
    "  ## 数据范围\n"
    "  [保留多久、刷新频率]\n"
    "  ## 下游\n"
    "  [谁在用这个数据]",
    size=11, color=SUB)

add_footer(slide)


# ═══ SLIDE 7: 实操 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "5️⃣  实操：给项目加文档")

add_code_block(slide, 400000, 1050000, 5000000, 2500000, """# 先看看你项目的文档
dbt docs generate
dbt docs serve

# 浏览器打开
open http://localhost:8080

# 停止服务
Ctrl+C""", size=13, color=ACCENT)

add_card(slide, 5600000, 1050000, 3200000, 2500000, CARD_BG)
add_bullets(slide, 5750000, 1150000, 2900000, 2300000, [
    "dbt docs generate：",
    ("读取 model + source 的描述", 1),
    ("从数据库读取字段类型", 1),
    ("生成 catalog.json + HTML", 1),
    ("", 0),
    "dbt docs serve：",
    ("启动本地 8080 端口", 1),
    ("实时查看文档站", 1),
], size=12, color=BODY)

add_card(slide, 400000, 3800000, 8300000, 2500000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 600000, 3900000, 7900000, 300000, "📊  你可以在文档站里看到：", size=16, color=WHITE, bold=True)

items = [
    "project 总览：11 个 model + 31 个测试 + 90 个 seed",
    "fct_orders 详情页：docs block 富文本描述 + 字段口径 + 测试列表",
    "血缘图：stg_sales_order_header → fct_orders 的可视化链路",
    "tag 筛选：core / finance 快速定位核心模型",
]
for i, item in enumerate(items):
    add_box(slide, 600000, 4250000 + i * 400000, 7900000, 350000, f"• {item}", size=12, color=BODY)

add_footer(slide)


# ═══ SLIDE 8: 同声传译 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "🛠️  同声传译 — 工程师视角")

concepts = [
    ("dbt docs generate", "Unity Catalog lineage + COMMENTS", "一键文档站，不手写 Word"),
    ("lineage graph", "Spark UI DAG 可视化", "数据血缘链路可视化"),
    ("schema.yml description", "COMMENT ON TABLE/COLUMN", "表和字段描述随代码管理"),
    ("docs block", "Markdown 长文档", "口径说明、使用文档"),
    ("dbt docs serve", "mkdocs / sphinx-autobuild", "本地 HTTP 预览"),
]

y_start = 1100000
for i, (dbt_concept, spark_concept, plain) in enumerate(concepts):
    y = y_start + i * 850000
    add_card(slide, 400000, y, 2500000, 700000, CARD_BG)
    add_box(slide, 450000, y + 50000, 2400000, 200000, dbt_concept, size=12, color=TITLE, bold=True)
    add_box(slide, 450000, y + 300000, 2400000, 350000, spark_concept, size=10, color=BODY)

    add_card(slide, 3000000, y, 2700000, 700000, CARD_BG)
    add_box(slide, 3050000, y + 50000, 2600000, 200000, "你熟悉的场景", size=10, color=ORANGE, bold=True)
    add_box(slide, 3050000, y + 300000, 2600000, 350000, spark_concept, size=10, color=BODY)

    add_card(slide, 5800000, y, 3000000, 700000, CARD_BG)
    add_box(slide, 5850000, y + 50000, 2900000, 200000, "大白话", size=10, color=ACCENT, bold=True)
    add_box(slide, 5850000, y + 300000, 2900000, 350000, plain, size=10, color=BODY)

add_footer(slide)


# ═══ SLIDE 9: 直觉陷阱 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "⚠️  工程师直觉陷阱")

add_card(slide, 500000, 1200000, 8100000, 1200000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 700000, 1300000, 7800000, 1000000,
    "🔥 直觉陷阱\n"
    '文档写在 yaml 里？不是有数据目录工具（Alation / Atlan / Collibra）吗？',
    size=15, color=RED)

add_card(slide, 500000, 2600000, 8100000, 1500000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 2700000, 7800000, 1300000,
    "✅ dbt 的哲学\n"
    "文档应该离数据最近——跟 model 放在同一个 PR 里。\n"
    "改代码时顺手更新文档，代码审查顺带审了。\n\n"
    "数据目录工具管的是\"数据资产目录\"\n"
    "dbt docs 管的是\"这张表怎么算出来的\"\n"
    "两者互补，不互斥。",
    size=14, color=ACCENT)

add_card(slide, 500000, 4300000, 8100000, 1200000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 700000, 4400000, 7800000, 1000000,
    "💡 实战映射\n"
    "药企平台里 Unity Catalog 管元数据治理、\n"
    "dbt docs 管口径说明——各司其职。\n"
    "dbt 负责血缘和业务口径，UC 负责治理和权限。",
    size=13, color=SUB)
add_footer(slide)


# ═══ SLIDE 10: 下集预告 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1500000, 7300000, 800000, "✅  本集完成", size=40, color=GREEN, bold=True)

add_box(slide, 900000, 2500000, 7300000, 300000, "你已学会：", size=16, color=WHITE, bold=True)
add_bullets(slide, 900000, 2800000, 7300000, 1500000, [
    "在 schema.yml 中写三层次描述（model/column/source）",
    "dbt docs generate + docs serve 一键生成文档站",
    "docs block 自定义富文本说明书",
    "tag 标签筛选核心模型",
    "文档最佳实践：什么该写 / 什么不该写",
], size=14, color=BODY)

add_divider(slide, 4500000)

add_box(slide, 900000, 4700000, 7300000, 400000, "下一集：P07 — 宏和包管理：告别重复代码", size=18, color=TITLE, bold=True)

add_box(slide, 900000, 5400000, 7300000, 400000, "配套代码：github.com/raycdut/dbt-learn", size=14, color=MUTED)
add_box(slide, 900000, 5700000, 7300000, 400000, "我是 Nick，下期见 👋", size=14, color=SUB)


# ── Save ──
output_path = "/Users/chendong/Projects/dbt-learn/ppt/P06-文档.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
