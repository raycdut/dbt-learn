#!/usr/bin/env python3
"""Rebuild P08-快照-SCD-Type-2.pptx with dark theme + full content."""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = 9144000
SLIDE_H = 6858000

BG      = RGBColor(0x1a, 0x1b, 0x26)
TITLE   = RGBColor(0x00, 0x9d, 0xff)
SUB     = RGBColor(0x88, 0xbb, 0xdd)
WHITE   = RGBColor(0xff, 0xff, 0xff)
BODY    = RGBColor(0xcc, 0xcc, 0xdd)
ACCENT  = RGBColor(0x00, 0xcc, 0x99)
ORANGE  = RGBColor(0xff, 0x99, 0x44)
MUTED   = RGBColor(0x66, 0x77, 0x88)
GREEN   = RGBColor(0x00, 0xcc, 0x66)
RED     = RGBColor(0xff, 0x55, 0x55)
YELLOW  = RGBColor(0xff, 0xcc, 0x44)
CARD_BG = RGBColor(0x22, 0x24, 0x32)
CODE_BG = RGBColor(0x2a, 0x2b, 0x3a)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_box(slide, left, top, w, h, text, size=18, color=BODY, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, w, h, items, size=14, color=BODY, title=None, title_color=TITLE):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = title_color
        p.font.bold = True
        first = False
    for item in items:
        txt, lvl = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.text = f"{'  ' * lvl}◆ {txt}"
        p.font.size = Pt(size - lvl * 2)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.level = lvl
        first = False
    return txBox

def add_card(slide, left, top, w, h, fill=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_code_block(slide, left, top, w, h, code_text, size=11, color=ACCENT):
    add_card(slide, left, top, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(Emu(left + 100000), Emu(top + 80000), Emu(w - 200000), Emu(h - 160000))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Courier New"

def add_header(slide, text, color=TITLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x15, 0x20)
    shape.line.fill.background()
    add_box(slide, 400000, 120000, 8300000, 600000, text, size=26, color=color, bold=True)

def add_footer(slide, text="P08 - 快照 (SCD Type 2) | Nick"):
    add_box(slide, 400000, SLIDE_H - 400000, 8000000, 300000, text, size=10, color=MUTED)

def add_divider(slide, y=3200000):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 400000, Emu(y), Emu(8300000), Emu(3000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
    shape.line.fill.background()


# ═══ SLIDE 1: Title ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1800000, 7300000, 1200000, "P08", size=72, color=TITLE, bold=True, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 2800000, 7300000, 600000, "快照 — 一行配置搞定历史数据", size=36, color=WHITE, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 3400000, 7300000, 400000, "dbt 从入门到精通", size=22, color=SUB, align=PP_ALIGN.LEFT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 900000, Emu(3900000), Emu(2000000), Emu(3000))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
add_box(slide, 900000, 4100000, 7300000, 400000, "Nick  |  github.com/raycdut/dbt-learn", size=16, color=MUTED)
add_box(slide, 900000, 4400000, 7300000, 400000, "《高级工程师视角下的 dbt 通关指南》系列", size=14, color=MUTED)


# ═══ SLIDE 2: 什么是快照 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "📋  什么是快照（Snapshot）？")

add_box(slide, 400000, 1050000, 8300000, 300000, "场景：客户改了地址，你想知道「改之前的地址是什么」", size=16, color=WHITE, bold=True)

# Timeline diagram as text
add_code_block(slide, 400000, 1500000, 8300000, 1200000, """时间线：
┌──────────────┬──────────────┬──────────────┐
│  2026-01-01  │  2026-03-15  │  2026-06-01  │
│  地址: 北京   │  地址: 上海   │  地址: 深圳   │
│  status: A   │  status: A   │  status: B   │
└──────────────┴──────────────┴──────────────┘""", size=11, color=YELLOW)

add_card(slide, 400000, 2900000, 4000000, 1800000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 500000, 3000000, 3800000, 250000, "普通表", size=20, color=RED, bold=True)
add_bullets(slide, 500000, 3300000, 3800000, 1300000, [
    "只保留最新状态（地址 = 深圳）",
    "丢失了全部历史",
    "查不到「去年客户地址是什么」",
], size=13, color=BODY)

add_card(slide, 4700000, 2900000, 4000000, 1800000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 4800000, 3000000, 3800000, 250000, "快照表", size=20, color=GREEN, bold=True)
add_bullets(slide, 4800000, 3300000, 3800000, 1300000, [
    "每次变化新增一条记录",
    "看到所有历史版本",
    "可查询任意时间点的「切片」",
], size=13, color=BODY)

add_card(slide, 400000, 4900000, 8300000, 1500000, RGBColor(0x14, 0x15, 0x20))
add_box(slide, 600000, 5000000, 7900000, 1200000,
    "💡 简单说：快照 = 带时间戳的版本控制，每改一次就留一个副本。\n"
    "普通表 UPDATE，快照表 INSERT + 旧行标记失效。",
    size=14, color=ACCENT)
add_footer(slide)


# ═══ SLIDE 3: dbt snapshot 做了什么 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "🔍  dbt snapshot 做了什么？")

add_box(slide, 400000, 1000000, 8300000, 300000, "db snapshot = 一行配置把普通表变成 SCD Type 2 版本表", size=14, color=SUB, bold=False)

add_code_block(slide, 400000, 1400000, 4200000, 2500000, """raw_customers（源表）
┌─────┬───────┬────────┐
│ id  │ name  │ city   │
├─────┼───────┼────────┤
│  1  │ 张三  │ 北京   │
│  1  │ 张三  │ 上海   │
│  1  │ 张三  │ 深圳   │
└─────┴───────┴────────┘

(源表每行是变更记录)""", size=9, color=YELLOW)

add_code_block(slide, 4700000, 1400000, 4200000, 2500000, """customers_snapshot（快照表）
┌─────┬───────┬────────┬────────────┬──────────┐
│ id  │ name  │ city   │ valid_from │ valid_to │
├─────┼───────┼────────┼────────────┼──────────┤
│  1  │ 张三  │ 北京   │ 2026-01-01 │ 2026-03-15│
│  1  │ 张三  │ 上海   │ 2026-03-15 │ 2026-06-01│
│  1  │ 张三  │ 深圳   │ 2026-06-01 │ NULL     │
└─────┴───────┴────────┴────────────┴──────────┘
         ↑ 额外字段：dbt 自动添加""", size=9, color=ACCENT)

add_card(slide, 400000, 4100000, 8300000, 2300000, RGBColor(0x14, 0x15, 0x20))
add_box(slide, 600000, 4200000, 7900000, 300000, "dbt snapshot 自动做的事情：", size=15, color=WHITE, bold=True)
add_bullets(slide, 600000, 4550000, 7900000, 1700000, [
    "比较源表当前数据和快照表最新版本",
    "发现变化 → 旧行设置 dbt_valid_to，新行 INSERT",
    "没有变化 → 什么都不做",
    "物理删除 → 设置 dbt_valid_to（如果 invalidate_hard_deletes=True）",
    "所有版本控制字段自动维护，零手写",
], size=12, color=BODY)
add_footer(slide)


# ═══ SLIDE 4: 配置语法 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "1️⃣  配置快照 — 基本语法")

add_code_block(slide, 400000, 1050000, 5000000, 2600000, """-- snapshots/customers_snapshot.sql
{% snapshot customers_snapshot %}

    {{
        config(
            target_schema='snapshots',
            unique_key='customer_id',
            strategy='timestamp',
            updated_at='updated_at',
            invalidate_hard_deletes=True,
        )
    }}

    SELECT * FROM {{ source('seed_data', 'raw_customers') }}

{% endsnapshot %}""", size=10, color=ACCENT)

add_card(slide, 5600000, 1050000, 3200000, 2600000, CARD_BG)
add_bullets(slide, 5750000, 1150000, 2900000, 2400000, [
    ("参数说明：", 0),
    ("target_schema", 0),
    ("快照存放的 schema", 1),
    ("unique_key", 0),
    ("业务主键（辨认为同一实体）", 1),
    ("strategy", 0),
    ("timestamp（推荐）/ check_cols", 1),
    ("updated_at", 0),
    ("源表的时间戳字段", 1),
    ("invalidate_hard_deletes", 0),
    ("源表物理删除时标记失效", 1),
    ("", 0),
    ("运行：", 0),
    ("dbt snapshot", 1),
    ("（跟 dbt run / seed 同级）", 1),
], size=11, color=BODY)

add_divider(slide, 3900000)

# Two strategies
strat_y = 4100000
add_card(slide, 400000, strat_y, 4100000, 2200000, CARD_BG)
add_box(slide, 500000, strat_y + 80000, 3900000, 250000, "timestamp 策略（推荐）", size=15, color=TITLE, bold=True)
add_code_block(slide, 500000, strat_y + 400000, 3800000, 600000, """strategy='timestamp'
updated_at='updated_at'
# 源表需要有更新时间戳字段""", size=11, color=ACCENT)
add_box(slide, 500000, strat_y + 1200000, 3900000, 800000,
    "比较 max(updated_at) 判断变更\n"
    "性能好，推荐首选\n"
    "源表必须有时间戳", size=11, color=BODY)

add_card(slide, 4600000, strat_y, 4100000, 2200000, CARD_BG)
add_box(slide, 4700000, strat_y + 80000, 3900000, 250000, "check_cols 策略（备选）", size=15, color=ORANGE, bold=True)
add_code_block(slide, 4700000, strat_y + 400000, 3800000, 600000, """strategy='check_cols'
check_cols=['city', 'email']
# 没有时间戳时比较指定列""", size=11, color=ACCENT)
add_box(slide, 4700000, strat_y + 1200000, 3900000, 800000,
    "比较指定列的当前 vs 历史值\n"
    "check_cols='all'（慎用，性能差）\n"
    "相当于做全字段 diff", size=11, color=BODY)
add_footer(slide)


# ═══ SLIDE 5: 快照表结构 + 查询 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "2️⃣  快照表的结构 + 查询方式")

# Table headers
cols_data = [
    ("字段", "作用", "示例"),
    ("dbt_scd_id", "唯一行标识（字段 hash）", "abc123..."),
    ("dbt_valid_from", "该版本开始时间", "2026-01-01"),
    ("dbt_valid_to", "该版本结束时间", "2026-03-15（当前为 NULL）"),
    ("dbt_updated_at", "最后更新时间", "2026-03-15"),
]
y = 1050000
for i, (col1, col2, col3) in enumerate(cols_data):
    bg_c = CARD_BG if i > 0 else RGBColor(0x14, 0x15, 0x20)
    add_card(slide, 400000, y, 2400000, 350000, bg_c)
    add_box(slide, 450000, y + 50000, 2300000, 250000, col1, size=11, color=TITLE if i == 0 else ACCENT, bold=(i == 0))
    add_card(slide, 2900000, y, 2600000, 350000, bg_c)
    add_box(slide, 2950000, y + 50000, 2500000, 250000, col2, size=11, color=TITLE if i == 0 else BODY, bold=(i == 0))
    add_card(slide, 5600000, y, 3200000, 350000, bg_c)
    add_box(slide, 5650000, y + 50000, 3100000, 250000, col3, size=11, color=TITLE if i == 0 else MUTED, bold=(i == 0))
    y += 370000

add_divider(slide, 2900000)

# Query examples
add_box(slide, 400000, 3050000, 8300000, 300000, "三种常用查询", size=18, color=WHITE, bold=True)

add_code_block(slide, 400000, 3400000, 8300000, 1000000, """-- 1. 查当前数据：dbt_valid_to IS NULL
SELECT * FROM snapshots.customers_snapshot WHERE dbt_valid_to IS NULL

-- 2. 查某个客户的历史：ORDER BY dbt_valid_from
SELECT * FROM snapshots.customers_snapshot WHERE customer_id = 101 ORDER BY dbt_valid_from

-- 3. 查时间点快照：时间 BETWEEN valid_from AND COALESCE(valid_to, '9999-12-31')
SELECT * FROM snapshots.customers_snapshot
WHERE '2026-02-01' BETWEEN dbt_valid_from AND COALESCE(dbt_valid_to, '9999-12-31')""", size=10, color=ACCENT)

add_card(slide, 400000, 4600000, 8300000, 1800000, RGBColor(0x14, 0x15, 0x20))
add_box(slide, 600000, 4700000, 7900000, 1500000,
    "💡 你的项目已配：\n"
    "• snapshots/customers_snapshot.sql — check_cols 策略，只监控 city 变化\n"
    "• seeds/raw_customers.csv — 3 个客户的演示数据\n"
    "• sources.yml → seed_data source 已就绪\n"
    "• 运行：dbt seed → dbt snapshot → 查询 snapshots 表",
    size=12, color=ACCENT)
add_footer(slide)


# ═══ SLIDE 6: 实操 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "3️⃣  实操：创建 + 运行快照")

# Step 1
add_box(slide, 400000, 1050000, 8300000, 300000, "Step 1：导入 seed + 首次运行 snapshot", size=15, color=TITLE, bold=True)
add_code_block(slide, 400000, 1400000, 5000000, 800000, """# 导入 seed 数据
dbt seed

# 首次 snapshot（此时每人只有 1 条记录）
dbt snapshot

# 查看
duckdb dev.duckdb -c \\
  "SELECT * FROM snapshots.customers_snapshot" """, size=11, color=ACCENT)

add_card(slide, 5600000, 1400000, 3200000, 800000, CARD_BG)
add_box(slide, 5750000, 1500000, 2900000, 600000,
    "首次 snapshot 会把全部数据\n"
    "作为第一个版本写入\n"
    "dbt_valid_to = NULL\n"
    "（数据还没变过）", size=11, color=BODY)

# Step 2
add_box(slide, 400000, 2400000, 8300000, 300000, "Step 2：模拟数据变化", size=15, color=TITLE, bold=True)
add_code_block(slide, 400000, 2750000, 5000000, 800000, """# 把 raw_customers.csv 中张三的城市
# 从 北京 改成 上海，重新导入

duckdb dev.duckdb -c \\
  "DROP TABLE IF EXISTS raw_customers;"
dbt seed

# 再次 snapshot
dbt snapshot""", size=11, color=ACCENT)

add_card(slide, 5600000, 2750000, 3200000, 800000, CARD_BG)
add_box(slide, 5750000, 2850000, 2900000, 600000,
    "dbt 发现 city 变了\n"
    "→ 旧行设 dbt_valid_to\n"
    "→ 新行 INSERT\n"
    "张三现在有 2 条记录", size=11, color=BODY)

# Step 3
add_box(slide, 400000, 3750000, 8300000, 300000, "Step 3：查看历史版本", size=15, color=TITLE, bold=True)
add_code_block(slide, 400000, 4100000, 8300000, 1200000, """duckdb dev.duckdb -c "
SELECT customer_id, name, city, dbt_valid_from, dbt_valid_to
FROM snapshots.customers_snapshot
ORDER BY customer_id, dbt_valid_from
"

# 结果应该是：
# 101 | 张三 | 北京 | 2026-... | 2026-...  (旧版本)
# 101 | 张三 | 上海 | 2026-... | NULL       (当前版本)
# 102 | 李四 | 上海 | 2026-... | NULL       (没变过)
# 103 | 王五 | 深圳 | 2026-... | NULL       (没变过)""", size=9, color=ACCENT)
add_footer(slide)


# ═══ SLIDE 7: invalidate_hard_deletes ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "4️⃣  invalidate_hard_deletes — 处理物理删除")

add_card(slide, 400000, 1050000, 8300000, 1500000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 600000, 1150000, 7900000, 1200000,
    "❌ 默认行为：源表物理删除了某行 → 快照表什么也不做\n"
    "  → 你的 snapshot 里这条记录一直显示为当前有效\n"
    "  → 但实际上数据已经不存在了\n\n"
    "✅ invalidate_hard_deletes=True\n"
    "  → 下一轮 snapshot 时，把 dbt_valid_to 设为当前时间\n"
    "  → 标记为「已删除」，不再有效",
    size=13, color=BODY)

add_divider(slide, 2800000)

add_card(slide, 400000, 3000000, 8300000, 1200000, CARD_BG)
add_box(slide, 600000, 3100000, 7900000, 1000000,
    "invalidate_hard_deletes 的原理：\n"
    "1. snapshot 查询源表的所有 unique_key\n"
    "2. 跟快照表中 dbt_valid_to IS NULL 的 unique_key 做差集\n"
    "3. 差集中的 key → 源表已删除 → 设置 dbt_valid_to",
    size=12, color=BODY)

add_card(slide, 400000, 4400000, 8300000, 2000000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 600000, 4500000, 7900000, 300000, "⚡ 快照最佳实践", size=16, color=WHITE, bold=True)
add_bullets(slide, 600000, 4900000, 7900000, 1400000, [
    ("✅ 适合快照：客户信息 / 产品分类 / 组织架构", 0),
    ("❌ 不适合：不可变日志 / 事实表 / 超大表", 0),
    ("快照清理：DELETE FROM snapshots... WHERE dbt_valid_to < 365 days", 0),
    ("查看膨胀率：GROUP BY customer_id ORDER BY COUNT(*) DESC", 0),
], size=12, color=BODY)
add_footer(slide)


# ═══ SLIDE 8: SCD Type 对比 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "📊  SCD Type 1 / 2 / 3 对比")

add_box(slide, 400000, 1000000, 8300000, 300000, "假设客户的地址从北京改成上海：", size=14, color=SUB)

add_code_block(slide, 400000, 1400000, 8300000, 1500000, """Type 1（覆盖）：         Type 2（保留历史）：      Type 3（保留上次）：
┌────────┬────────┐    ┌────────┬────────┬──────┐  ┌────────┬────────┬────────┐
│ city   │        │    │ city   │ valid_ │ valid│  │ city   │ prev_  │ prev_  │
│        │        │    │        │ from   │ to   │  │        │ city   │ date   │
├────────┼────────┤    ├────────┼────────┼──────┤  ├────────┼────────┼────────┤
│ 上海   │ ← 覆盖  │    │ 北京   │ 01-01  │ 03-15│  │ 上海   │ 北京   │ 03-15  │
└────────┴────────┘    │ 上海   │ 03-15  │ NULL │  └────────┴────────┴────────┘
  北京丢了              └────────┴────────┴──────┘    只保留了上一次
                            都在""", size=10, color=YELLOW)

# Comparison table
rows = [
    ("维度", "Type 1（覆盖）", "Type 2（dbt snapshot）", "Type 3（手写）"),
    ("历史保留", "❌ 全丢", "✅ 全保留", "⚠️ 只保留上次"),
    ("数据量", "一行", "一行→N行", "一行+2字段"),
    ("查当前", "SELECT *", "WHERE valid_to IS NULL", "SELECT *"),
    ("查历史", "❌ 无法", "✅ 任意时间点", "⚠️ 只能看上一步"),
    ("实现难度", "UPDATE", "{% snapshot %}一行配置", "手写 MERGE"),
    ("适用场景", "不关心历史", "审计/合规/时间旅行", "快速回看"),
]

y = 3000000
for i, (c1, c2, c3, c4) in enumerate(rows):
    bg_c = CARD_BG if i > 0 else RGBColor(0x14, 0x15, 0x20)
    w1, w2, w3, w4 = 1800000, 2100000, 2100000, 2100000
    add_card(slide, 400000, y, w1, 300000, bg_c)
    add_box(slide, 450000, y + 50000, w1 - 100000, 200000, c1, size=10, color=TITLE if i == 0 else WHITE, bold=(i == 0))
    add_card(slide, 400000 + w1, y, w2, 300000, bg_c)
    add_box(slide, 450000 + w1, y + 50000, w2 - 100000, 200000, c2, size=10, color=TITLE if i == 0 else BODY, bold=(i == 0))
    add_card(slide, 400000 + w1 + w2, y, w3, 300000, bg_c)
    add_box(slide, 450000 + w1 + w2, y + 50000, w3 - 100000, 200000, c3, size=10, color=TITLE if i == 0 else ACCENT, bold=(i == 0))
    add_card(slide, 400000 + w1 + w2 + w3, y, w4, 300000, bg_c)
    add_box(slide, 450000 + w1 + w2 + w3, y + 50000, w4 - 100000, 200000, c4, size=10, color=TITLE if i == 0 else ORANGE, bold=(i == 0))
    y += 310000

add_card(slide, 400000, 5300000, 8300000, 1100000, RGBColor(0x14, 0x15, 0x20))
add_box(slide, 600000, 5400000, 7900000, 900000,
    "生产选型：90% → Type 2（dbt snapshot），8% → Type 1（默认 merge），2% → Type 3（手写 LAG）\n"
    "Type 2 一年膨胀 3-5x 行数，百万客户 ≈ 500 万行/年，对 Snowflake 没压力",
    size=13, color=SUB)
add_footer(slide)


# ═══ SLIDE 9: 直觉陷阱 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "⚠️  工程师直觉陷阱")

add_card(slide, 500000, 1200000, 8100000, 1500000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 700000, 1300000, 7800000, 1300000,
    "🔥 直觉陷阱 1\n"
    "Type 2 数据膨胀太吓人了，感觉 Type 1 就够了\n\n"
    "✅ 现实\n"
    "典型企业维度表一年变化 3-5 次。100 万客户 × 5 次 = 500 万行/年。\n"
    "对 Snowflake 来说这点量跟没有一样。\n"
    "真正要小心的是每天全量 snapshot — 用 timestamp 策略只在下游变化时跑。",
    size=13, color=BODY)

add_card(slide, 500000, 2900000, 8100000, 1500000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 3000000, 7800000, 1300000,
    "🔥 直觉陷阱 2\n"
    "Type 3 面试经常被问到，应该很常用\n\n"
    "✅ 现实\n"
    "Type 3 面试题多，实战极少。\n"
    "只保留上一次的边界太尴尬了 — 查两次变动它就丢了，审计追溯完全不达标。\n"
    "90% 场景 Type 2 够用，剩下 10% 不需要历史记录。",
    size=13, color=BODY)

add_card(slide, 500000, 4600000, 8100000, 1800000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 700000, 4700000, 7800000, 1600000,
    "💡 实战映射 — 药企场景\n"
    "客户维度信息经常变（部门调整、地址变更、联系人换人）。\n"
    "用 snapshot 记录历史变更，做审计追溯时直接查：\n"
    '"去年这个时候 CRM 里客户地址存的是什么？"\n\n'
    "SELECT city FROM snapshots.customers_snapshot\n"
    "WHERE customer_id = 101 AND '2025-06-01' BETWEEN dbt_valid_from AND COALESCE(dbt_valid_to, '9999-12-31')",
    size=12, color=ACCENT)
add_footer(slide)


# ═══ SLIDE 10: 同声传译 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "🛠️  同声传译 — 工程师视角")

concepts = [
    ("SCD Type 2", "维度历史变化追踪", "客户地址变了，新旧都留着"),
    ("dbt snapshot", "Delta CDC + 时间旅行", "一键开 SCD Type 2，零手写 MERGE"),
    ("timestamp strategy", "比较 max(updated_at)", "按时间戳判断行有没有更新"),
    ("check_cols strategy", "exceptAll 比较字段值", "没时间戳的时候比字段值"),
    ("dbt_valid_from/to", "Delta __start/__end_date", "版本有效期的起止时间"),
    ("invalidate_hard_deletes", "软删除标记", "源删了行，快照标记失效"),
]

y_start = 1100000
for i, (dbt_concept, spark_concept, plain) in enumerate(concepts):
    y = y_start + i * 750000
    add_card(slide, 400000, y, 2500000, 650000, CARD_BG)
    add_box(slide, 450000, y + 50000, 2400000, 180000, dbt_concept, size=12, color=TITLE, bold=True)
    add_box(slide, 450000, y + 280000, 2400000, 320000, spark_concept, size=10, color=BODY)

    add_card(slide, 3000000, y, 2700000, 650000, CARD_BG)
    add_box(slide, 3050000, y + 50000, 2600000, 180000, "你熟悉的场景", size=10, color=ORANGE, bold=True)
    add_box(slide, 3050000, y + 280000, 2600000, 320000, spark_concept, size=10, color=BODY)

    add_card(slide, 5800000, y, 3000000, 650000, CARD_BG)
    add_box(slide, 5850000, y + 50000, 2900000, 180000, "大白话", size=10, color=ACCENT, bold=True)
    add_box(slide, 5850000, y + 280000, 2900000, 320000, plain, size=10, color=BODY)

add_footer(slide)


# ═══ SLIDE 11: 总结 ════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1500000, 7300000, 800000, "✅  本集完成", size=40, color=GREEN, bold=True)

add_box(slide, 900000, 2500000, 7300000, 300000, "你已学会：", size=16, color=WHITE, bold=True)
add_bullets(slide, 900000, 2800000, 7300000, 1500000, [
    "什么是快照（带版本控制的维度表）",
    "dbt snapshot 的一行配置（timestamp / check_cols）",
    "快照自动添加的字段和三种常用查询",
    "SCD Type 1/2/3 对比与生产选型",
    "invalidate_hard_deletes 处理物理删除",
], size=14, color=BODY)

add_divider(slide, 4500000)

add_box(slide, 900000, 4700000, 7300000, 400000, "下一集：P09 — Hook + Operation：dbt 的生产化配置", size=18, color=TITLE, bold=True)

add_box(slide, 900000, 5400000, 7300000, 400000, "配套代码：github.com/raycdut/dbt-learn", size=14, color=MUTED)
add_box(slide, 900000, 5700000, 7300000, 400000, "我是 Nick，下期见 👋", size=14, color=SUB)


# ── Save ──
output_path = "/Users/chendong/Projects/dbt-learn/ppt/P08-快照-SCD-Type-2.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
