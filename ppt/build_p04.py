#!/usr/bin/env python3
"""Rebuild P04-增量模型.pptx with dark theme + full content."""

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──
SLIDE_W = 9144000   # 10 inches
SLIDE_H = 6858000   # 7.5 inches

# Color palette (dark theme)
BG      = RGBColor(0x1a, 0x1b, 0x26)
TITLE   = RGBColor(0x00, 0x9d, 0xff)
SUB     = RGBColor(0x88, 0xbb, 0xdd)
WHITE   = RGBColor(0xff, 0xff, 0xff)
BODY    = RGBColor(0xcc, 0xcc, 0xdd)
ACCENT  = RGBColor(0x00, 0xcc, 0x99)
ORANGE  = RGBColor(0xff, 0x99, 0x44)
MUTED   = RGBColor(0x66, 0x77, 0x88)
GREEN   = RGBColor(0x00, 0xcc, 0x66)
RED     = RGBColor(0xff, 0x55, 0x55)
YELLOW  = RGBColor(0xff, 0xcc, 0x44)
CARD_BG = RGBColor(0x22, 0x24, 0x32)
CODE_BG = RGBColor(0x2a, 0x2b, 0x3a)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank


# ── Helpers ──
def add_bg(slide, color=BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_box(slide, left, top, w, h, text, size=18, color=BODY, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, w, h, items, size=14, color=BODY, title=None, title_color=TITLE):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Arial"
        first = False
    for item in items:
        if isinstance(item, str):
            txt, lvl = item, 0
        else:
            txt, lvl = item
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.text = f"{'  ' * lvl}◆ {txt}"
        p.font.size = Pt(size - lvl * 2)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.level = lvl
        first = False
    return txBox

def add_card(slide, left, top, w, h, fill=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_code_block(slide, left, top, w, h, code_text, size=11, color=ACCENT):
    add_card(slide, left, top, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(Emu(left + 100000), Emu(top + 80000), Emu(w - 200000), Emu(h - 160000))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Courier New"

def add_header(slide, text, color=TITLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x15, 0x20)
    shape.line.fill.background()
    add_box(slide, 400000, 120000, 8300000, 600000, text, size=26, color=color, bold=True)

def add_footer(slide, text="P04 - 增量模型与物化策略 | Nick"):
    add_box(slide, 400000, SLIDE_H - 400000, 8000000, 300000, text, size=10, color=MUTED)

def add_divider(slide, y=3200000):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 400000, Emu(y), Emu(8300000), Emu(3000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
    shape.line.fill.background()


# ════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1800000, 7300000, 1200000, "P04", size=72, color=TITLE, bold=True, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 2800000, 7300000, 600000, "增量模型与物化策略", size=36, color=WHITE, bold=False, align=PP_ALIGN.LEFT)
add_box(slide, 900000, 3400000, 7300000, 400000, "dbt 从入门到精通", size=22, color=SUB, align=PP_ALIGN.LEFT)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 900000, Emu(3900000), Emu(2000000), Emu(3000))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_box(slide, 900000, 4100000, 7300000, 400000, "Nick  |  github.com/raycdut/dbt-learn", size=16, color=MUTED)
add_box(slide, 900000, 4400000, 7300000, 400000, "《高级工程师视角下的 dbt 通关指南》系列", size=14, color=MUTED)


# ════════════════════════════════════════════
# SLIDE 2: 为什么需要增量模型？
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "📋  为什么需要增量模型？")

add_card(slide, 500000, 1100000, 8100000, 1400000, CARD_BG)
add_box(slide, 700000, 1200000, 7800000, 300000, "场景：你有 3 亿行订单数据，每天新增 50 万行", size=16, color=WHITE, bold=True)
add_box(slide, 700000, 1550000, 7800000, 800000,
    "❌ full refresh 每天：删表重建，重读 3 亿行 → 40 分钟\n"
    "✅ incremental 每天：只读新增 50 万行 → 2 分钟",
    size=15, color=BODY)

add_card(slide, 500000, 2800000, 8100000, 1000000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 2900000, 7800000, 800000,
    "💡 核心思路：不做重复工作，只处理变化的数据",
    size=18, color=ACCENT, bold=True)

add_box(slide, 500000, 4100000, 8100000, 600000,
    "什么时候该换 incremental？\n"
    "当你感觉到\"dbt run 太慢了\"的时候——百万级以下全量就行，以上必须增量",
    size=13, color=SUB)

add_card(slide, 500000, 4800000, 8100000, 1200000, RGBColor(0x2a, 0x1a, 0x2a))
add_box(slide, 700000, 4900000, 7800000, 1000000,
    "⚡ 和 Delta Lake 的 MERGE 是同一件事\n"
    "你在 Databricks 写 MERGE INTO + 增量过滤 → dbt 自动帮你生成",
    size=14, color=ORANGE)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 3: 四种物化方式对比
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "1️⃣  四种物化方式对比")

# Table header
cols_x = [400000, 550000, 2200000, 3900000, 5600000, 7300000]
col_w = [150000, 1550000, 1600000, 1600000, 1600000, 1400000]
headers = ["", "怎么存的", "刷新方式", "适用场景", "数据量"]
for i, h in enumerate(headers):
    add_box(slide, cols_x[i], 1050000, col_w[i], 300000, h, size=11, color=TITLE, bold=True)

div_y = 1400000
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 400000, Emu(div_y), Emu(8300000), Emu(2000))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
shape.line.fill.background()

rows = [
    ("view", "SQL 定义", "每次查询实时算", "轻量过滤/stg层", "小"),
    ("table", "物理表", "每次 dbt run 重建", "小表/需要性能", "< 100万行"),
    ("incremental", "物理表", "只处理增量数据", "大表/每日增量", "任意"),
    ("ephemeral", "CTE 片段", "不落表/内联", "中间逻辑复用", "极小"),
]
for i, (name, store, refresh, scene, vol) in enumerate(rows):
    y = 1550000 + i * 700000
    add_card(slide, 400000, y, 8300000, 550000, CARD_BG)
    colors = [TITLE, BODY, BODY, BODY, BODY]
    texts = [name, store, refresh, scene, vol]
    for j, (tx, clr) in enumerate(zip(texts, colors)):
        add_box(slide, cols_x[j+1], y + 100000, col_w[j+1], 350000, tx, size=12, color=clr,
                bold=(j == 0))

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 4: 物化方式配置 & 2.2 配置方式
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "2️⃣  物化配置方式")

add_box(slide, 500000, 1100000, 4000000, 300000, "全局设置（dbt_project.yml）", size=16, color=WHITE, bold=True)
add_code_block(slide, 500000, 1450000, 4000000, 1600000, """models:
  my_first_project:
    staging:
      +materialized: view
    intermediate:
      +materialized: table
    marts:
      +materialized: incremental""", size=10, color=YELLOW)

add_box(slide, 500000, 3300000, 4000000, 300000, "单 model 设置（SQL 文件头）", size=16, color=WHITE, bold=True)
add_code_block(slide, 500000, 3650000, 4000000, 1200000, """{{ config(
    materialized='incremental',
    unique_key='order_id'
) }}""", size=11, color=ACCENT)

# Right side: explanation
add_bullets(slide, 4800000, 1100000, 4000000, 5000000, [
    "全局配置写在 dbt_project.yml",
    "按目录层级覆盖：staging / int / marts",
    "",
    "单 model 配置写在 SQL 头部",
    "config() 覆盖全局设置",
    "",
    "优先级：",
    ("model 级 config > 目录级 > 全局默认", 1),
    "",
    "💡 和 Spark 一样：",
    "全局 conf → session conf → SQL hint",
], size=13, color=BODY)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 5: Incremental 核心语法
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "3️⃣  Incremental 核心语法")

add_code_block(slide, 400000, 1100000, 5000000, 2500000, """{{ config(
    materialized='incremental',
    unique_key='order_id'
) }}

SELECT * FROM {{ ref('stg_orders') }}

{% if is_incremental() %}
  WHERE created_at >
    (SELECT MAX(created_at)
      FROM {{ this }})
{% endif %}""", size=11, color=ACCENT)

# Explanation panel
add_bullets(slide, 5600000, 1100000, 3200000, 5000000, [
    (""),  # padding
    ("第一次 dbt run：", 0),
    ("表不存在 → is_incremental = False", 1),
    ("走全量 SELECT，没有 WHERE", 1),
    ("创建目标表", 1),
    (""),
    ("第二次 dbt run：", 0),
    ("表已存在 → is_incremental = True", 1),
    ("WHERE 生效，只处理新数据", 1),
    ("执行 MERGE INTO + 增量过滤", 1),
    (""),
    ("dbt run --full-refresh：", 0),
    ("强制 is_incremental = False", 1),
    ("全量重建", 1),
], size=12, color=BODY)

# Key elements
add_card(slide, 400000, 3800000, 8300000, 2500000, CARD_BG)
add_box(slide, 600000, 3900000, 7900000, 250000, "✨ 关键语法", size=16, color=TITLE, bold=True)

elems = [
    ("is_incremental()", "判断目标表是否存在？"),
    ("{{ this }}", "指向目标表自身的引用"),
    ("unique_key", "MERGE 的 ON 条件，去重键"),
    ("incremental_strategy", "merge / delete+insert / append"),
]
for i, (kw, desc) in enumerate(elems):
    y = 4200000 + i * 500000
    add_box(slide, 600000, y, 2000000, 350000, kw, size=13, color=ORANGE, bold=True)
    add_box(slide, 2700000, y, 5500000, 350000, f"→ {desc}", size=12, color=BODY)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 6: 三种增量策略
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "4️⃣  三种增量策略")

# merge
add_card(slide, 400000, 1100000, 2700000, 2500000, CARD_BG)
add_box(slide, 500000, 1150000, 2500000, 300000, "① merge（默认）", size=15, color=TITLE, bold=True)
add_code_block(slide, 450000, 1500000, 2600000, 700000, """MERGE INTO t USING s
ON t.id = s.id
WHEN MATCHED UPDATE
WHEN NOT MATCHED INSERT""", size=8, color=ACCENT)
add_box(slide, 500000, 2300000, 2500000, 400000, "有则更新，无则插入\n适合：可更新的订单表", size=11, color=BODY)

# delete+insert
add_card(slide, 3200000, 1100000, 2700000, 2500000, CARD_BG)
add_box(slide, 3300000, 1150000, 2500000, 300000, "② delete+insert", size=15, color=ORANGE, bold=True)
add_code_block(slide, 3250000, 1500000, 2600000, 700000, """DELETE FROM t
WHERE id IN (SELECT id FROM s);
INSERT INTO t SELECT * FROM s;""", size=8, color=ACCENT)
add_box(slide, 3300000, 2300000, 2500000, 400000, "先删后插，适合：审计轨迹", size=11, color=BODY)

# append
add_card(slide, 6000000, 1100000, 2700000, 2500000, CARD_BG)
add_box(slide, 6100000, 1150000, 2500000, 300000, "③ append", size=15, color=GREEN, bold=True)
add_code_block(slide, 6050000, 1500000, 2600000, 700000, """INSERT INTO t
SELECT * FROM s;
-- 不检查 unique_key""", size=9, color=ACCENT)
add_box(slide, 6100000, 2300000, 2500000, 400000, "只追加不更新，最快\n适合：日志/事件流水", size=11, color=BODY)

# Strategy guide at bottom
add_card(slide, 400000, 3800000, 8300000, 2500000, CARD_BG)
add_box(slide, 600000, 3900000, 7900000, 250000, "📊  策略选择指南", size=16, color=WHITE, bold=True)

guide = [
    ("订单表（可更新）", "merge", "支持 upsert"),
    ("订单表（只增不改）", "append", "更快"),
    ("快照维度表", "merge", "需要覆盖"),
    ("日志表", "append", "只增不改"),
    ("需要审计轨迹", "delete+insert", "先删再插"),
]
for i, (scene, strategy, reason) in enumerate(guide):
    y = 4200000 + i * 420000
    add_box(slide, 600000, y, 2500000, 350000, scene, size=11, color=BODY)
    add_box(slide, 3200000, y, 1500000, 350000, strategy, size=12, color=ACCENT, bold=True)
    add_box(slide, 4800000, y, 3000000, 350000, f"→ {reason}", size=11, color=MUTED)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 7: 实战改造 — 项目里做了什么
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "5️⃣  实战改造 — fct_orders 增量")

add_code_block(slide, 400000, 1100000, 5000000, 3200000, """{{ config(
    materialized='incremental',
    unique_key='sales_order_id',
    incremental_strategy='merge'
) }}

SELECT
    h.sales_order_id,
    h.order_date,
    h.customer_id,
    h.total_due,
    h.status_code,
    CASE h.status_code ...
      END AS status_label
FROM {{ ref('stg_sales_order_header') }} h
LEFT JOIN {{ ref('stg_customer') }} c
  ON h.customer_id = c.customer_id

{% if is_incremental() %}
  WHERE h.order_date >
    (SELECT MAX(order_date)
      FROM {{ this }})
{% endif %}""", size=9, color=ACCENT)

# Results
add_card(slide, 5600000, 1100000, 3200000, 2700000, CARD_BG)
add_box(slide, 5750000, 1200000, 2900000, 300000, "📈 运行结果", size=16, color=GREEN, bold=True)

results = [
    ("1st run", "Full load", "31,465 行"),
    ("2nd run", "Incremental", "0 行（无新数据）"),
    ("--full-refresh", "Recreate", "31,465 行重建"),
]
for i, (run, mode, rows) in enumerate(results):
    y = 1600000 + i * 600000
    add_box(slide, 5800000, y, 900000, 300000, run, size=11, color=ORANGE, bold=True)
    add_box(slide, 6800000, y, 900000, 300000, mode, size=11, color=BODY)
    add_box(slide, 7800000, y, 900000, 300000, rows, size=11, color=GREEN)

add_box(slide, 5600000, 3500000, 3200000, 300000,
    "→ 新增文件：", size=14, color=WHITE, bold=True)
add_box(slide, 5600000, 3800000, 3200000, 2500000,
    "models/marts/fct_orders.sql\n"
    "models/marts/_marts__models.yml\n\n"
    "测试（4 个）：\n"
    "unique(sales_order_id)\n"
    "not_null(sales_order_id)\n"
    "not_null(order_date)\n"
    "not_null(customer_id)",
    size=12, color=BODY)

# Airflow notice
add_card(slide, 400000, 4500000, 8300000, 1800000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 600000, 4600000, 7900000, 300000, "☁️  生产部署注意", size=16, color=TITLE, bold=True)
add_box(slide, 600000, 4950000, 7900000, 1200000,
    "• 日常调度：dbt run（自动增量 MERGE）\n"
    "• 全量重建：dbt run --full-refresh 或 Airflow 传参 {\"full_refresh\": true}\n"
    "• Cosmos 集成：DbtDag + DuckDBLocalProfileMapping，无需 profiles.yml",
    size=12, color=BODY)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 8: 常用命令速查
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "📌  常用命令速查")

cmds = [
    ("dbt run", "增量/全量按配置执行"),
    ("dbt run --full-refresh", "强制全量重建所有 model"),
    ("dbt run --select fct_orders", "只跑指定 model"),
    ("dbt run --select fct_orders --full-refresh", "只重建指定 model"),
    ("dbt build", "seed + run + test 一步到位"),
    ("dbt test --select fct_orders", "只跑特定 model 的测试"),
]
for i, (cmd, desc) in enumerate(cmds):
    y = 1100000 + i * 850000
    add_card(slide, 500000, y, 3700000, 650000, CODE_BG)
    add_box(slide, 600000, y + 100000, 3500000, 450000, cmd, size=14, color=YELLOW, bold=True)
    add_box(slide, 4400000, y + 100000, 4300000, 450000, desc, size=14, color=BODY)

# Tip box
add_card(slide, 500000, 6200000, 8100000, 1, RGBColor(0x33, 0x44, 0x55))

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 9: 同声传译
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "🛠️  同声传译 — 工程师视角")

concepts = [
    ("materialized='view'", "Spark SQL CREATE VIEW", "每次查询实时算"),
    ("materialized='table'", "Spark SQL CTAS", "先算好存起来，查表"),
    ("materialized='incremental'", "Delta MERGE INTO", "只处理增量，不全量重跑"),
    ("materialized='ephemeral'", "WITH CTE 子句", "中间表不落盘"),
    ("is_incremental()", "检查目标 Delta 表存在", "第一次全量，之后增量"),
    ("incremental_strategy='merge'", "Delta MERGE ON key", "有则更新无则插入"),
    ("on_schema_change", "Spark schema evolution", "源表加列了怎么处理"),
]

y_start = 1100000
for i, (dbt_concept, spark_concept, plain) in enumerate(concepts):
    y = y_start + i * 700000
    add_card(slide, 400000, y, 2500000, 550000, CARD_BG)
    add_box(slide, 450000, y + 50000, 2400000, 150000, dbt_concept, size=12, color=TITLE, bold=True)
    add_box(slide, 450000, y + 200000, 2400000, 300000, spark_concept, size=10, color=BODY)

    add_card(slide, 3000000, y, 2700000, 550000, CARD_BG)
    add_box(slide, 3050000, y + 50000, 2600000, 150000, "你熟悉的场景", size=10, color=ORANGE, bold=True)
    add_box(slide, 3050000, y + 200000, 2600000, 300000, spark_concept, size=10, color=BODY)

    add_card(slide, 5800000, y, 3000000, 550000, CARD_BG)
    add_box(slide, 5850000, y + 50000, 2900000, 150000, "大白话", size=10, color=ACCENT, bold=True)
    add_box(slide, 5850000, y + 200000, 2900000, 300000, plain, size=10, color=BODY)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 10: 直觉陷阱
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "⚠️  工程师直觉陷阱")

add_card(slide, 500000, 1200000, 8100000, 1200000, RGBColor(0x2a, 0x1a, 0x1a))
add_box(slide, 700000, 1300000, 7800000, 1000000,
    "🔥 直觉陷阱\n"
    "\"增量模型好复杂，全量建 view 不是更简单吗？\"",
    size=16, color=RED)

add_card(slide, 500000, 2600000, 8100000, 1500000, RGBColor(0x1a, 0x2a, 0x1a))
add_box(slide, 700000, 2700000, 7800000, 1300000,
    "✅ 事实\n"
    "• view 简单但每次查询都重新算，百万级数据慢到怀疑人生\n"
    "• incremental 第一次麻烦，之后每次只处理增量\n"
    "• 就像 Delta 表做 MERGE INTO — 第一次全量写入，之后每天只写增量\n"
    "• dbt 用配置声明，连 MERGE SQL 都不用你手写",
    size=14, color=ACCENT)

add_card(slide, 500000, 4300000, 8100000, 1200000, RGBColor(0x1a, 0x1a, 0x2a))
add_box(slide, 700000, 4400000, 7800000, 1000000,
    "💡 实战映射\n"
    "在 A 股量化系统里，每日行情数据就是完美的增量场景——\n"
    "T 日只处理 T 日的数据，不需要重算历史。\n"
    "用 dbt 的 incremental mode 配合 DuckDB 的增量策略，\n"
    "比手写 Python 脚本干净得多。",
    size=13, color=SUB)
add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 11: 常见问题
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header(slide, "❓  常见问题")

problems = [
    ("第一次跑增量 model 很慢？", "正常。第一次是全量构建，第二次开始才是增量。"),
    ("增量跑错了，想全量重建？", "dbt run --full-refresh --select model_name"),
    ("怎么确认增量生效了？", "看 log：看到 \"incremental model\" 就是增量模式"),
    ("unique_key 选错会怎样？", "重复数据。选业务主键或 surrogate key。"),
    ("on_schema_change 怎么配？", "源表加列时用 sync_all_columns 自动同步结构"),
]

for i, (q, a) in enumerate(problems):
    y = 1200000 + i * 1050000
    add_box(slide, 600000, y, 8000000, 300000, f"Q：{q}", size=15, color=YELLOW, bold=True)
    add_box(slide, 800000, y + 400000, 7800000, 500000, f"A：{a}", size=13, color=BODY)
    if i < len(problems) - 1:
        add_divider(slide, y + 900000)

add_footer(slide)


# ════════════════════════════════════════════
# SLIDE 12: 下集预告
# ════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(60000))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE
shape.line.fill.background()

add_box(slide, 900000, 1500000, 7300000, 800000, "✅  本集完成", size=40, color=GREEN, bold=True)

add_box(slide, 900000, 2500000, 7300000, 300000, "你已学会：", size=16, color=WHITE, bold=True)
add_bullets(slide, 900000, 2800000, 7300000, 1500000, [
    "四种物化方式的取舍：view / table / incremental / ephemeral",
    "incremental 核心语法：is_incremental / this / unique_key",
    "三种增量策略：merge / delete+insert / append",
    "实战：fct_orders 改增量 + marts 测试定义",
    "Airflow Cosmos 集成模式",
], size=14, color=BODY)

add_divider(slide, 4500000)

add_box(slide, 900000, 4700000, 7300000, 400000, "下一集：P05 — 测试与数据质量", size=18, color=TITLE, bold=True)

add_box(slide, 900000, 5400000, 7300000, 400000, "配套代码：github.com/raycdut/dbt-learn", size=14, color=MUTED)
add_box(slide, 900000, 5700000, 7300000, 400000, "我是 Nick，下期见 👋", size=14, color=SUB)


# ── Save ──
output_path = "/Users/chendong/Projects/dbt-learn/ppt/P04-增量模型.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
